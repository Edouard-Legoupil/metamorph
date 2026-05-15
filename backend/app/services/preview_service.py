"""
File Preview Service

Handles generation of previews for various file types discovered during website crawling.
Supports PDF, Word, Excel, PowerPoint, text, HTML, and other common document formats.
"""

import os
import tempfile
import hashlib
from typing import Optional, Dict, Any
from pathlib import Path

import requests
from fastapi import HTTPException

# Import file processing libraries
try:
    import PyPDF2
    import pythoncom
    import win32com.client as win32
    from bs4 import BeautifulSoup
    from docx import Document as DocxDocument
    from pptx import Presentation
    from openpyxl import load_workbook
    from markdown import markdown
    import json
    import xml.etree.ElementTree as ET
    HAS_ALL_DEPENDENCIES = True
except ImportError as e:
    HAS_ALL_DEPENDENCIES = False
    print(f"Warning: Some preview dependencies missing: {e}")

from app.models.sql.website import DiscoveredFile, FileType
from app.core.config import settings


class PreviewService:
    """Service for generating file previews"""
    
    def __init__(self):
        self.cache_dir = Path(settings.preview_cache_dir or "/tmp/preview_cache")
        self.cache_dir.mkdir(exist_ok=True)
        self.max_preview_length = settings.max_preview_length or 10000
        self.max_file_size = settings.max_preview_file_size or 10 * 1024 * 1024  # 10MB
    
    def _get_cache_key(self, file_url: str, content_hash: Optional[str] = None) -> str:
        """Generate cache key for file preview"""
        if content_hash:
            return content_hash
        return hashlib.md5(file_url.encode('utf-8')).hexdigest()
    
    def _get_cached_preview(self, cache_key: str) -> Optional[str]:
        """Check if preview is cached"""
        cache_file = self.cache_dir / f"{cache_key}.preview"
        if cache_file.exists():
            return cache_file.read_text()
        return None
    
    def _cache_preview(self, cache_key: str, preview_text: str) -> None:
        """Cache preview for future use"""
        cache_file = self.cache_dir / f"{cache_key}.preview"
        cache_file.write_text(preview_text)
    
    def _download_file(self, url: str) -> tempfile.NamedTemporaryFile:
        """Download file from URL for preview generation"""
        try:
            response = requests.get(url, timeout=30, stream=True)
            response.raise_for_status()
            
            # Check file size
            content_length = int(response.headers.get('content-length', 0))
            if content_length > self.max_file_size:
                raise HTTPException(
                    status_code=413,
                    detail=f"File too large for preview ({content_length} bytes > {self.max_file_size} bytes)"
                )
            
            # Create temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{url.split('.')[-1]}")
            for chunk in response.iter_content(chunk_size=8192):
                if chunk:
                    temp_file.write(chunk)
            temp_file.close()
            return temp_file
            
        except requests.RequestException as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to download file for preview: {str(e)}"
            )
    
    def _cleanup_temp_file(self, temp_file_path: str) -> None:
        """Clean up temporary file"""
        try:
            os.unlink(temp_file_path)
        except OSError:
            pass
    
    def _truncate_text(self, text: str, max_length: int = 1000) -> str:
        """Truncate text to maximum length"""
        if len(text) <= max_length:
            return text
        return text[:max_length] + "... (truncated)"
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return self._truncate_text(text)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from PDF: {str(e)}"
            )
    
    def _extract_text_from_word(self, file_path: str) -> str:
        """Extract text from Word document (DOCX)"""
        try:
            doc = DocxDocument(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return self._truncate_text(text)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from Word document: {str(e)}"
            )
    
    def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from Excel spreadsheet"""
        try:
            wb = load_workbook(filename=file_path, read_only=True)
            text = ""
            for sheet in wb.sheetnames:
                text += f"\n=== Sheet: {sheet} ===\n"
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            return self._truncate_text(text)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from Excel: {str(e)}"
            )
    
    def _extract_text_from_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                text += f"\n=== Slide {slide.slide_id} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return self._truncate_text(text)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from PowerPoint: {str(e)}"
            )
    
    def _extract_text_from_html(self, file_path: str) -> str:
        """Extract text from HTML file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text()
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                return self._truncate_text(text)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from HTML: {str(e)}"
            )
    
    def _extract_text_from_text(self, file_path: str) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                return self._truncate_text(text)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to read text file: {str(e)}"
            )
    
    def _extract_text_from_markdown(self, file_path: str) -> str:
        """Extract text from Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                # Convert markdown to plain text by removing formatting
                html = markdown(text)
                soup = BeautifulSoup(html, 'html.parser')
                return self._truncate_text(soup.get_text())
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from Markdown: {str(e)}"
            )
    
    def _extract_text_from_json(self, file_path: str) -> str:
        """Extract text from JSON file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                return self._truncate_text(json.dumps(data, indent=2))
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from JSON: {str(e)}"
            )
    
    def _extract_text_from_xml(self, file_path: str) -> str:
        """Extract text from XML file"""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            text = ET.tostring(root, encoding='unicode')
            return self._truncate_text(text)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to extract text from XML: {str(e)}"
            )
    
    def generate_preview(self, discovered_file: DiscoveredFile) -> Dict[str, Any]:
        """Generate preview for a discovered file"""
        # Check cache first
        cache_key = self._get_cache_key(discovered_file.url, discovered_file.content_hash)
        cached_preview = self._get_cached_preview(cache_key)
        if cached_preview:
            return {
                "file_id": discovered_file.id,
                "url": discovered_file.url,
                "file_name": discovered_file.file_name,
                "file_type": discovered_file.file_type.value,
                "preview": cached_preview,
                "cached": True
            }
        
        # Download file
        temp_file = None
        try:
            temp_file = self._download_file(discovered_file.url)
            
            # Generate preview based on file type
            preview_text = ""
            if discovered_file.file_type == FileType.PDF:
                preview_text = self._extract_text_from_pdf(temp_file.name)
            elif discovered_file.file_type in [FileType.WORD, FileType.DOCX]:
                preview_text = self._extract_text_from_word(temp_file.name)
            elif discovered_file.file_type in [FileType.EXCEL, FileType.XLSX]:
                preview_text = self._extract_text_from_excel(temp_file.name)
            elif discovered_file.file_type in [FileType.POWERPOINT, FileType.PPTX]:
                preview_text = self._extract_text_from_powerpoint(temp_file.name)
            elif discovered_file.file_type in [FileType.HTML, FileType.HTM]:
                preview_text = self._extract_text_from_html(temp_file.name)
            elif discovered_file.file_type == FileType.TEXT:
                preview_text = self._extract_text_from_text(temp_file.name)
            elif discovered_file.file_type == FileType.MARKDOWN:
                preview_text = self._extract_text_from_markdown(temp_file.name)
            elif discovered_file.file_type == FileType.JSON:
                preview_text = self._extract_text_from_json(temp_file.name)
            elif discovered_file.file_type == FileType.XML:
                preview_text = self._extract_text_from_xml(temp_file.name)
            else:
                preview_text = f"Preview not supported for file type: {discovered_file.file_type.value}"
            
            # Cache the preview
            self._cache_preview(cache_key, preview_text)
            
            return {
                "file_id": discovered_file.id,
                "url": discovered_file.url,
                "file_name": discovered_file.file_name,
                "file_type": discovered_file.file_type.value,
                "preview": preview_text,
                "cached": False
            }
            
        finally:
            if temp_file:
                self._cleanup_temp_file(temp_file.name)


# Singleton instance
preview_service = PreviewService()